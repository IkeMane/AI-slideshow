from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from openai import OpenAI
import base64
from PIL import Image
import io
import json
import os
from dotenv import load_dotenv

# For creating a PowerPoint presentation
from pptx import Presentation
from pptx.util import Inches


load_dotenv()
api_key = os.getenv('OPENAI_API_KEY')
Client = OpenAI(api_key=api_key)


def get_image_as_base64(file_path):
    """
    Encodes an image as a base64 string.

    :param file_path: Path to the image file.
    :return: Base64-encoded string of the image.
    """
    with Image.open(file_path) as image:
        buffered = io.BytesIO()
        image.save(buffered, format="PNG")  # You can change the format as needed
        img_str = base64.b64encode(buffered.getvalue()).decode('utf-8')
    return img_str



def create_image_data(file_path):
    """
    Creates a dictionary object that includes the base64-encoded image.
    
    :param file_path: Path to the image file.
    :return: Dictionary with the image data.
    """
    base64_image = get_image_as_base64(file_path)
    return {
        "image_url": {
            "content": base64_image,
            "type": "base64"
        }
    }




def create_presentation(json_data):
    prs = Presentation()
    for slide_data in json_data['presentation']['slides']:
        slide = add_slide(prs)
        if slide_data['type'] == 'table':
            add_table(slide, slide_data['title'], slide_data['headers'], slide_data['rows'])
        elif slide_data['type'] == 'text':
            add_text_slide(slide, slide_data['title'], slide_data['content'])
    return prs



def add_slide(presentation, layout=5):
    slide_layout = presentation.slide_layouts[layout]
    return presentation.slides.add_slide(slide_layout)


def add_table(slide, title, headers, rows):
    # Set slide background color to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # RGB value for white

    # Add and format the title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(18)  # Adjust the font size as needed

    # Add the table
    table = slide.shapes.add_table(1 + len(rows), len(headers), Inches(1), Inches(1.5), Inches(8), Inches(0.8)).table
    table.autofit = False  # Disable autofit so we can set fixed column widths

    # Set header styles
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(137, 163, 211)  # RGB value for blue
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.name = "Arial"
            paragraph.font.size = Pt(13)  # Adjust the font size as needed
            paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text

    # Set row styles
    for row_index, row in enumerate(rows, start=1):
        for col_index, cell_text in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)  # Adjust the font size as needed
                paragraph.font.name = "Arial"

    # Set column widths
    num_cols = len(headers)
    table_width = Inches(8)
    col_width = table_width / num_cols
    for col_index in range(num_cols):
        for row_index in range(1 + len(rows)):
            cell = table.cell(row_index, col_index)
            cell.width = col_width



def add_text_slide(slide, title, content):
    # Set the title text and font size
    title_shape = slide.shapes.title 
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(16)  # Set the title font size to 18pt
    
    # Add the textbox for content
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    # Add each line of content
    for line in content:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(13)  # Set the content font size if needed





# Function to encode the image
def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')


def build_messages(base64_image):
    return{
            "role": "user",
            "content": [
                    {
                    "type": "text",
                    "text": "Here's and example format for the JSON youre required to provide based off the given image be sure to match the headers with the table size. just leave it with a space or something:"
                    },
                    {
                    "type": "text",
                    "text": """{
                        "presentation": {
                            "title": "title of the presentation",
                            "slides": [                   
                            {
                                "type": "table",
                                "title": "slide title",
                                "headers": ["header", "header..."],
                                "rows": [
                                ["column content", "column content..."],
                                ["column content", "column content..."],
                                ["column content", "column content..."]
                                ]
                            },
                            ]
                        }
                        }"""
                    },
                    {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{base64_image}"
                        }
                    }
            ]
        }


# Example usage:
if __name__ == "__main__":
    directory = '/Users/isaacharmon/Files/Code Stuff/accounting-class-help/AI-slideshow/screenshots'
    presentation = Presentation()
    messages = list()
    messages.append({"role": "system", "content": "I read the image and provide the content in a consistant structured json format. I can create text or table slides based off the image im given. I try to answer the question in the image if possable. I only create one slide per question. I try to keep the slides simple and easy to read for the students. Show the answers."})     
    for filename in os.listdir(directory):
        if filename.endswith('.png'):
            file_path = os.path.join(directory, filename)
            base64_image = encode_image(file_path)


            messages.append(build_messages(base64_image)) #add image message to conversation

            response = Client.chat.completions.create(
                model="gpt-4o",
                messages=messages,
                response_format={ "type": "json_object" },
            )
    
            message = response.choices[0].message
            message_text = message.content
            messages.append({"role": "assistant", "content": message_text}) # add to conversation

            print(f"Response: {message_text}")

            #JSON read and processing
            json_data = message_text

            slide_data = json.loads(json_data)
            for slide_info in slide_data['presentation']['slides']:
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
                if slide_info['type'] == 'table':
                    add_table(slide, slide_info['title'], slide_info['headers'], slide_info['rows'])
                elif slide_info['type'] == 'text':
                    add_text_slide(slide, slide_info['title'], slide_info['content'])


    presentation.save('AI_Presentation.pptx')
    print("Presentation created successfully!")
# ````````